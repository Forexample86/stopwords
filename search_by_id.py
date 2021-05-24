import csv
import subprocess
import os
import logging
import warnings as w
from collections import Counter

import docx
from pptx import Presentation
from scripts._blacklist import BlackList
from scripts._work import Work
import codecs
import gitlab


def searcher_rep(project_ssh):
    """
    Подключение к gitwork и выгрузка проекта по ssh
    :param project_ssh:
    :return:
    """
    # git@gitwork.ru:polev/test.git
    # https://gitwork.ru/polev/test.git

    try:
        gl = BlackList()
        gl = gl.get_connect()
    except gitlab.GitlabAuthenticationError as err:
        w.warn(err.error_message, Warning)
        return err.error_message

    # Получаем проект
    get = get_project(project_ssh)

    # fullname
    try:
        projects = gl.projects.list()
        for project in projects:
            if project.ssh_url_to_repo == project_ssh:
                project_ = project
                print(project_.namespace["name"])
                return project_.namespace["name"]
    except gitlab.GitlabListError as err:
        w.warn(err.error_message, Warning)


def searcher_files():
    """
    :return: paths
    """
    suffix = ('.docx', '.doc', '.pptx', '.md')
    paths = []
    folder = os.getcwd()
    for root, dirs, files in os.walk(folder):
        for file in files:
            if file.endswith(suffix) and not file.startswith('~'):
                paths.append(os.path.join(root, file))
    #  Конвертер doc в docx(работает в linux через libre office)
    #  args = ['soffice', '--headless', '--convert-to', 'docx', ]
    #  for count in paths:
    #      if count.endswith('.doc'):
    #          print(count)
    #          subprocess.Popen([args, count], stdout=subprocess.PIPE)
    return paths


def search_expansion():
    """
    :return: расширения
    """
    suffix = ('.docx', '.doc', '.pptx', '.ppt', '.md', '.txt', '.rtf')
    expansion = []
    folder = os.getcwd()
    for root, dirs, files in os.walk(folder):
        for file in files:
            if file.endswith(suffix[0:7]) and not file.startswith('~'):
                expansion.append(file.split(".")[-1])
    return expansion


def get_project(project_ssh):
    # выкачиваем проект
    args = ['git', 'clone', project_ssh]
    res = subprocess.Popen(args, stdout=subprocess.PIPE)
    out, error = res.communicate()
    if not error:
        print(out)
        return out
    else:
        print(error)
        return error


def csv_out(data, path):
    """
    Выводит в файл csv
    :param: file_obj
    :return:
    """
    dict_data = [Counter(data)]
    print(dict_data)
    csv_columns = ['docx', 'txt', 'doc', 'md', 'pptx', 'rtf']
    with open(path, "w") as csv_file:
        writer = csv.DictWriter(csv_file, fieldnames=csv_columns, dialect='excel',
                                extrasaction='raise', restval='', delimiter=':')
        writer.writeheader()
        for row in dict_data:
            writer.writerow(row)


def black_list(file):
    with codecs.open(file, 'r', encoding='utf-8') as fin:
        black = fin.read().splitlines()
    return black


def output(fullname, word, file):
    text = f"fullname:  {fullname}  ' | '  stop word: {word}"
    with codecs.open(file, "a", encoding='utf-8') as fin:
        out = fin.write(text + '\n')
        fin.close()
    return out


def pars_files(paths, fullname, file_b, file_o):
    print("Pars .md")
    work = Work(paths, fullname, file_b, file_o)

    for elem in paths:
        if elem.endswith('.md' or '.txt'):
            check_word(elem, work)

        if elem.endswith('.pptx' or '.ppt'):
            check_pptx(elem, work)

        if elem.endswith('.docx'):
            check_docx(elem, work)


def check_word(elem, work):
    with codecs.open(elem, encoding='utf-8') as openfile:
        for line in openfile:
            for key in black_list(work.file_b):
                if key in line:
                    print(f"I found {key} word")
                    output(work.fullname, key, work.file_o)
                else:
                    # print("Words not found!")
                    continue
                return key


def check_pptx(elem, work):
    prs = Presentation(elem)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.lower()
                for key in black_list(work.file_b):
                    if key in shape.text:
                        print(f"I found {key} word")
                        output(work.fullname, key, work.file_o)
                    else:
                        # print("Words not found!")
                        continue
                    return key


def check_docx(elem, work):
    doc = docx.Document(elem)
    for paragraph in doc.paragraphs:
        text = paragraph.text.lower()
        for key in black_list(work.file_b):
            if key in text:
                print(f"I found {key} word")
                output(work.fullname, key, work.file_o)
            else:
                # print("Words not found!")
                continue
            return key

#
# def pars_docx(paths, fullname, file_b, file_o):
#     print("Pars .docx")
#
#     for elem in paths:
#         if elem.endswith('.docx'):
#             doc = docx.Document(elem)
#             for paragraph in doc.paragraphs:
#                 text = paragraph.text.lower()
#                 for key in black_list(file_b):
#                     if key in text:
#                         print(f"I found {key} word")
#                         output(fullname, key, file_o)
#                     else:
#                         # print("Words not found!")
#                         continue
#                     return key
#
#
# def pars_pptx(paths, fullname, file_b, file_o):
#     print("Pars .pptx")
#
#     for elem in paths:
#         if elem.endswith('.pptx' or '.ppt'):
#             prs = Presentation(elem)
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         shape.text = shape.text.lower()
#                         for key in black_list(file_b):
#                             if key in shape.text:
#                                 print(f"I found {key} word")
#                                 output(fullname, key, file_o)
#                             else:
#                                 # print("Words not found!")
#                                 continue
#                             return key