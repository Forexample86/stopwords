"""
Функции для работы
"""
import os

import docx
from pptx import Presentation
from scripts._blacklist import BlackList
import logging
import warnings as w
import codecs
import csv
import subprocess

from scripts._work import Work
from collections import Counter


def black_list(file):
    with codecs.open(file, 'r', encoding='utf-8') as fin:
        black = fin.read().splitlines()
    return black


def pars_files(paths, fullname, file_b, file_o):
    # fle_b - blacklist file
    # paths -
    """
    Функция парсинга файлов
    Args:
        paths: пути
        fullname:
        file_b:
        file_o:

    Returns:

    """
    print("Pars .md")
    work = Work(paths, fullname, file_b, file_o)
    for elem in paths:
        if elem.endswith('.md' or '.txt'):
            check_word(elem, work)

        if elem.endswith('.pptx' or '.ppt'):
            check_pptx(elem, work)

        if elem.endswith('.docx'):
            check_docx(elem, work)


def check_word(elem, work):
    """
    Проверка word
    Args:
        elem: изучаемый объект
        work: класс с переменными

    Returns:

    """
    with codecs.open(elem, encoding='utf-8') as openfile:
        for line in openfile:
            for key in black_list(work.file_b):
                if key in line:
                    print(f"I found {key} word")
                    output(work, key)
                else:
                    # print("Words not found!")
                    continue
                return key


def check_pptx(elem, work):
    """
    Проверка pptx
    Args:
        elem: изучаемый объект
        work: класс с переменными

    Returns:

    """
    prs = Presentation(elem)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.lower()
                for key in black_list(work.file_b):
                    if key in shape.text:
                        print(f"I found {key} word")
                        output(work, key)
                    else:
                        # print("Words not found!")
                        continue
                    return key


def check_docx(elem, work):
    """
    Проверка docx
    Args:
        elem: изучаемый объект
        work: класс с переменными

    Returns:

    """
    doc = docx.Document(elem)
    for paragraph in doc.paragraphs:
        text = paragraph.text.lower()
        for key in black_list(work.file_b):
            if key in text:
                print(f"I found {key} word")
                output(work, key)
            else:
                # print("Words not found!")
                continue
            return key


def return_paths():
    """
    Функция возврата файлов с расширениями
    :return: paths
    """
    suffix = ('.docx', '.doc', '.pptx', '.md')
    paths = []
    folder = os.getcwd()
    for root, dirs, files in os.walk(folder):
        for file in files:
            if file.endswith(suffix) and not file.startswith('~'):
                paths.append(os.path.join(root, file))
    #  Конвертер doc в docx(работает в linux через libre office)
    #  args = ['soffice', '--headless', '--convert-to', 'docx', ]
    #  for count in paths:
    #      if count.endswith('.doc'):
    #          print(count)
    #          subprocess.Popen([args, count], stdout=subprocess.PIPE)
    return paths


def search_expansion():
    """
    :return: расширения
    """
    suffix = ('.docx', '.doc', '.pptx', '.ppt', '.md', '.txt', '.rtf')
    expansion = []
    folder = os.getcwd()
    for root, dirs, files in os.walk(folder):
        for file in files:
            if file.endswith(suffix[0:7]) and not file.startswith('~'):
                expansion.append(file.split(".")[-1])
    return expansion


def output(work, key):
    text = f"fullname:  {work.fullname}  ' | '  stop word: {key}"
    with codecs.open(work.file_o, "a", encoding='utf-8') as fin:
        out = fin.write(text + '\n')
        fin.close()
    return out


def csv_out(data, path):
    """
    Выводит в файл csv
    :param: file_obj - название файла

    :return:
    """
    dict_data = [Counter(data)]
    print(dict_data)
    csv_columns = ['docx', 'txt', 'doc', 'md', 'pptx', 'rtf']
    with open(path, "w") as csv_file:
        writer = csv.DictWriter(csv_file, fieldnames=csv_columns, dialect='excel',
                                extrasaction='raise', restval='', delimiter=':')
        writer.writeheader()
        for row in dict_data:
            writer.writerow(row)


def get_project(project_ssh):
    # выкачиваем проект
    args = ['git', 'clone', project_ssh]
    res = subprocess.Popen(args, stdout=subprocess.PIPE)
    out, error = res.communicate()
    if not error:
        print(out)
        return out
    else:
        print(error)
        return error