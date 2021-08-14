"""
Функции для работы
"""
import os
import re
import shutil
import stat
import codecs
import csv

import subprocess
from os import path
import pandas as pd

import docx
import pymorphy2
from pptx import Presentation

from scripts._work import Work


def black_list(file):
    """
    Получение списка плохих слов
    :param file: откуда брать
    :return: список
    """
    with codecs.open(file, 'r', encoding='utf-8') as fin:
        black = fin.read().splitlines()
    return black


def pars_files(paths, fullname, file_b, file_o, file_csv):
    """
    Обработка парсинга
    :param paths:  путь
    :param fullname: имя пользователя
    :param file_b: файл с плохими словами
    :param file_o: файл для логгирования
    :param file_csv: файл для статистики
    """

    work = Work(paths, fullname, file_b, file_o, file_csv)
    for elem in paths:
        if elem.endswith('.md' or '.txt'):
            check_word(elem, work)

        if elem.endswith('.pptx' or '.ppt'):
            check_pptx(elem, work)

        if elem.endswith('.docx'):
            check_docx(elem, work)


def check_word(elem, work):
    """
    Проверка word
    Args:
        elem: изучаемый объект
        work: класс с переменными

    Returns:

    """
    black = black_list(work.file_b)
    with codecs.open(elem, encoding='utf-8') as openfile:
        for line in openfile:
            key = parse(line, black)
            if key:
                output(work, key)
                csv_out('word', work.file_csv)


def check_pptx(elem, work):
    """
    Проверка pptx
    Args:
        elem: изучаемый объект
        work: класс с переменными

    Returns:

    """
    prs = Presentation(elem)
    black = black_list(work.file_b)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.lower()
                key = parse(shape.text, black)
                if key:
                    output(work, key)
                    csv_out('pptx', work.file_csv)


def check_docx(elem, work):
    """
    Проверка docx
    Args:
        elem: изучаемый объект
        work: класс с переменными

    Returns:

    """
    doc = docx.Document(elem)
    black = black_list(work.file_b)

    for paragraph in doc.paragraphs:
        text = paragraph.text.lower()
        key = parse(text, black)
        if key:
            output(work, key)
            csv_out('docx', work.file_csv)


def parse(text, blacklist):
    """
    Парс файлов
    :param text: текст для парса
    :param blacklist: список плохих слов
    :return:
    """
    morph = pymorphy2.MorphAnalyzer()
    for key in blacklist:
        word = morph.parse(key)[0]
        if 'NOUN' in word.tag.POS:
            for i in ['nomn', 'gent', 'datv', 'accs', 'ablt', 'loct']:
                search = r"\b" + word.inflect({i}).word + r"\S*\b"
                if re.findall(search, text.lower()):
                    print(f"Найдено слово {key}")
                    return key

            for i in ['nomn', 'gent', 'datv', 'accs', 'ablt', 'loct']:
                search = r"\b" + word.inflect({i, 'plur'}).word + r"\S*\b"
                if re.findall(search, text.lower()):
                    print(f"Найдено слово {key}")
                    return key


def return_paths():
    """
    Функция возврата файлов с расширениями
    :return: paths
    """
    suffix = ('.docx', '.doc', '.pptx', '.md')
    paths = []
    folder = os.getcwd()
    for root, dirs, files in os.walk(folder):
        for file in files:
            if file.endswith(suffix) and not file.startswith('~'):
                paths.append(os.path.join(root, file))
    #  Конвертер doc в docx(работает в linux через libre office)
    #  args = ['soffice', '--headless', '--convert-to', 'docx', ]
    #  for count in paths:
    #      if count.endswith('.doc'):
    #          print(count)
    #          subprocess.Popen([args, count], stdout=subprocess.PIPE)
    return paths


def search_expansion():
    """
    Поиск необходимых расширений
    :return: расширения
    """
    suffix = ('.docx', '.doc', '.pptx', '.ppt', '.md', '.txt', '.rtf')
    expansion = []
    folder = os.getcwd()
    for root, dirs, files in os.walk(folder):
        for file in files:
            if file.endswith(suffix[0:7]) and not file.startswith('~'):
                expansion.append(file.split(".")[-1])
    return expansion


def output(work, key):
    """
    Вывод в файл пользователя и найденное плохое слово
    :param work: класс для работы
    :param key: плохое слово
    :return:
    """
    text = f"fullname:  {work.fullname}  ' | '  stop word: {key}"
    with codecs.open(work.file_o, "a", encoding='utf-8') as fin:
        out = fin.write(text + '\n')
        fin.close()
    return out


def csv_out(data, path_to_csv):
    """
    Выводит в файл csv
    :param: data - нужная колонка
    :param: path_to_csv - путь до файла

    :return:
    """
    """
    dict_data = [Counter(data)]
    print(dict_data)
    csv_columns = ['docx', 'txt', 'doc', 'md', 'pptx', 'rtf']
    with open(path_to_csv, "w") as csv_file:
        writer = csv.DictWriter(csv_file, fieldnames=csv_columns, dialect='excel',
                                extrasaction='raise', restval='', delimiter=':')
        writer.writeheader()
        for row in dict_data:
            writer.writerow(row)

    """
    df = pd.read_csv(path_to_csv, sep=':')
    df.loc[0, data] += 1

    df.to_csv(path_to_csv, index=False, sep=':')


def get_project(project_ssh):
    """
    :param project_ssh:
    :return: Вывод
    """
    # выкачиваем проект
    args = ['git', 'clone', project_ssh]
    """
    res = subprocess.Popen(args, stdout=subprocess.PIPE)
    out, error = res.communicate()
    if not error:
        return out
    print(error)
    return error
    """
    try:
        subprocess.check_output(args)
    except subprocess.CalledProcessError as e:
        raise OSError('Ошибка загрузки') from e


def get_proj_name(ssh):
    """
    Выделение имя проекта из ssh
    :param ssh: строка
    :return: имя проекта
    """
    name = ssh.rpartition('/')[2]
    return name.rpartition('.')[0]


def delete_project(name):
    """
    Удаление проекта
    :param name: название
    :return: None
    """
    try:
        path_to_dir = './' + name
        for root, dirs, files in os.walk(path_to_dir):
            for dir in dirs:
                os.chmod(path.join(root, dir), stat.S_IRWXU)
            for file in files:
                os.chmod(path.join(root, file), stat.S_IRWXU)
        shutil.rmtree(path_to_dir)
    except FileNotFoundError as f:
        raise FileNotFoundError(f.strerror + ' ' + f.filename) from f


def clear_file(file):
    """
    Очистка файла логгирования перед запуском
    :return: None
    """
    with open(file, "w") as f:
        f.truncate(0)
        f.close()
    return True
